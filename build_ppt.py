"""Build a 5-slide deck for SAGGY."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
OUT = HERE / "SAGGY_overview.pptx"
SAGGY_HEAD = HERE / "head saggy.png"
A3_LOGO = HERE / "a3brands-1.png"

# Palette pulled from index.html
NAVY = RGBColor(0x0A, 0x0F, 0x1F)
INK = RGBColor(0x12, 0x18, 0x2B)
LIME = RGBColor(0xC8, 0xFF, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xB8, 0xBE, 0xCC)
DIM = RGBColor(0x6B, 0x73, 0x86)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=WHITE, font="Calibri",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_eyebrow(slide, left, top, text):
    return add_text(slide, left, top, Inches(8), Inches(0.3), text,
                    size=11, bold=True, color=LIME, font="Consolas")


def add_pill(slide, left, top, text, *, fill=LIME, fg=NAVY):
    pad_w = Inches(0.18)
    # rough width based on text length
    width = Inches(max(1.2, 0.13 * len(text) + 0.4))
    height = Inches(0.42)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    tf = shape.text_frame
    tf.margin_left = pad_w
    tf.margin_right = pad_w
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Consolas"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = fg
    return shape


def add_card(slide, left, top, width, height, *, fill=INK, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_footer(slide):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "A3 BRANDS · GALAXY VISIBILITY DIAGNOSTIC",
             size=9, color=DIM, font="Consolas")
    add_text(slide, Inches(7), Inches(7.05), Inches(6), Inches(0.3),
             "saggy score · seo · aeo · geo",
             size=9, color=DIM, font="Consolas", align=PP_ALIGN.RIGHT)


# ---------- Slide 1: Title ----------
s1 = prs.slides.add_slide(BLANK)
set_bg(s1, NAVY)

# decorative lime bar
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.7), Inches(0.08), Inches(1.2))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = LIME

add_eyebrow(s1, Inches(0.75), Inches(0.75), "GALAXY · VISIBILITY DIAGNOSTIC")
add_text(s1, Inches(0.75), Inches(1.05), Inches(8.5), Inches(0.5),
         "Meet SAGGY.", size=20, bold=True, color=WHITE, font="Calibri")
add_text(s1, Inches(0.5), Inches(2.3), Inches(8.5), Inches(2),
         "Get your\nSAGGY Score.", size=64, bold=True, color=WHITE, font="Calibri")
# highlight word
hl = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.85), Inches(3.5), Inches(0.18))
hl.line.fill.background()
hl.fill.solid()
hl.fill.fore_color.rgb = LIME

add_text(s1, Inches(0.5), Inches(4.5), Inches(8.5), Inches(1.2),
         "How visible is your dealership across SEO, AEO,\nand Generative AI search?",
         size=20, color=GRAY)

# right-side stat pills
add_pill(s1, Inches(9.6), Inches(2.5), "10 QUESTIONS")
add_pill(s1, Inches(9.6), Inches(3.1), "UNDER 90 SECONDS")
add_pill(s1, Inches(9.6), Inches(3.7), "ONE SCORE / 100")

if SAGGY_HEAD.exists():
    s1.shapes.add_picture(str(SAGGY_HEAD), Inches(9.4), Inches(4.3), height=Inches(2.6))

add_footer(s1)


# ---------- Slide 2: The Problem ----------
s2 = prs.slides.add_slide(BLANK)
set_bg(s2, NAVY)

add_eyebrow(s2, Inches(0.5), Inches(0.6), "THE PROBLEM")
add_text(s2, Inches(0.5), Inches(0.95), Inches(12), Inches(1),
         "Your dealership is invisible where buyers\nactually search now.",
         size=34, bold=True, color=WHITE)

# three stat cards
cards = [
    ("AI search", "Buyers ask ChatGPT, Gemini, and Perplexity\nbefore they ever open Google."),
    ("Hidden leakage", "Leads slip out the back door of your\nfunnel — without your provider noticing."),
    ("Blind spend", "Without close rate by source, every\nbudget call is a guess."),
]
x = Inches(0.5)
y = Inches(3.0)
w = Inches(4.0)
h = Inches(3.3)
gap = Inches(0.2)
for i, (head, body) in enumerate(cards):
    cx = Emu(x + (w + gap) * i)
    add_card(s2, cx, y, w, h, fill=INK, border=DIM)
    add_text(s2, Emu(cx + Inches(0.3)), Emu(y + Inches(0.35)), Emu(w - Inches(0.6)), Inches(0.5),
             head.upper(), size=11, bold=True, color=LIME, font="Consolas")
    add_text(s2, Emu(cx + Inches(0.3)), Emu(y + Inches(0.95)), Emu(w - Inches(0.6)), Inches(2.3),
             body, size=16, color=WHITE)

add_footer(s2)


# ---------- Slide 3: How It Works ----------
s3 = prs.slides.add_slide(BLANK)
set_bg(s3, NAVY)
add_eyebrow(s3, Inches(0.5), Inches(0.6), "HOW IT WORKS")
add_text(s3, Inches(0.5), Inches(0.95), Inches(12), Inches(1),
         "Three pillars. One score. 90 seconds.",
         size=34, bold=True, color=WHITE)

pillars = [
    ("SEO", "Search Engine\nOptimization",
     "Classic ranking signals — schema, content depth, site architecture, crawlability."),
    ("AEO", "Answer Engine\nOptimization",
     "Snippet eligibility, structured Q&A, and review signals that win the answer box."),
    ("GEO", "Generative\nAI Search",
     "How often LLMs cite your inventory and brand when buyers ask in natural language."),
]
x = Inches(0.5)
y = Inches(2.4)
w = Inches(4.0)
h = Inches(4.0)
gap = Inches(0.2)
for i, (code, label, body) in enumerate(pillars):
    cx = Emu(x + (w + gap) * i)
    add_card(s3, cx, y, w, h, fill=INK)
    # giant pillar code
    add_text(s3, Emu(cx + Inches(0.3)), Emu(y + Inches(0.3)), Emu(w - Inches(0.6)), Inches(1.4),
             code, size=64, bold=True, color=LIME)
    add_text(s3, Emu(cx + Inches(0.3)), Emu(y + Inches(1.7)), Emu(w - Inches(0.6)), Inches(0.9),
             label, size=15, bold=True, color=WHITE)
    add_text(s3, Emu(cx + Inches(0.3)), Emu(y + Inches(2.55)), Emu(w - Inches(0.6)), Inches(1.4),
             body, size=13, color=GRAY)
    # /20 footer
    add_text(s3, Emu(cx + Inches(0.3)), Emu(y + Inches(3.55)), Emu(w - Inches(0.6)), Inches(0.3),
             "scored out of 20", size=10, color=DIM, font="Consolas")

add_footer(s3)


# ---------- Slide 4: What You Get ----------
s4 = prs.slides.add_slide(BLANK)
set_bg(s4, NAVY)
add_eyebrow(s4, Inches(0.5), Inches(0.6), "WHAT YOU GET")
add_text(s4, Inches(0.5), Inches(0.95), Inches(12), Inches(1),
         "A score, a verdict, and the top 3 fixes.",
         size=34, bold=True, color=WHITE)

# Big score block on left
add_card(s4, Inches(0.5), Inches(2.4), Inches(5.5), Inches(4.0), fill=INK, border=LIME)
add_text(s4, Inches(0.8), Inches(2.6), Inches(5), Inches(0.5),
         "SAGGY SCORE", size=12, bold=True, color=LIME, font="Consolas")
add_text(s4, Inches(0.8), Inches(3.0), Inches(5), Inches(2),
         "78 / 100", size=72, bold=True, color=WHITE)
add_text(s4, Inches(0.8), Inches(4.9), Inches(5), Inches(0.5),
         "Orbit Builder", size=20, bold=True, color=LIME)
add_text(s4, Inches(0.8), Inches(5.4), Inches(5), Inches(1),
         "“Solid orbit. The fundamentals are there.\nLeads are still slipping out the back door.”",
         size=12, color=GRAY)

# Right column: bands list
add_text(s4, Inches(6.4), Inches(2.4), Inches(6.5), Inches(0.4),
         "FIVE RESULT BANDS", size=12, bold=True, color=LIME, font="Consolas")
bands = [
    ("90+", "Lunar Leader", "Flag in the ground."),
    ("75+", "Orbit Builder", "Fundamentals locked."),
    ("60+", "Liftoff Stage", "Visible in search, invisible in AI."),
    ("40+", "Pre-Launch", "Engines not fired yet."),
    ("0+",  "Grounded", "Not in the game yet."),
]
row_y = Inches(2.85)
for i, (rng, name, blurb) in enumerate(bands):
    ry = Emu(row_y + Inches(0.65) * i)
    add_text(s4, Inches(6.4), ry, Inches(0.9), Inches(0.5),
             rng, size=16, bold=True, color=LIME, font="Consolas")
    add_text(s4, Inches(7.3), ry, Inches(2.4), Inches(0.5),
             name, size=14, bold=True, color=WHITE)
    add_text(s4, Inches(9.7), ry, Inches(3.5), Inches(0.5),
             blurb, size=12, color=GRAY)

add_footer(s4)


# ---------- Slide 5: CTA ----------
s5 = prs.slides.add_slide(BLANK)
set_bg(s5, NAVY)

# huge centered headline
add_eyebrow(s5, Inches(0.5), Inches(1.4), "READY?")
add_text(s5, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.2),
         "Find out exactly where you're leaking leads.",
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s5, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.6),
         "10 questions in Saggy’s voice. Under 90 seconds. One score out of 100.",
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

# CTA button (lime pill)
btn = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(4.4), Inches(4.4), Inches(4.5), Inches(0.85))
btn.adjustments[0] = 0.5
btn.line.fill.background()
btn.fill.solid()
btn.fill.fore_color.rgb = LIME
btf = btn.text_frame
btf.margin_top = Inches(0.1)
btf.margin_bottom = Inches(0.1)
p = btf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Book My Strategy Call →"
r.font.name = "Calibri"
r.font.size = Pt(20)
r.font.bold = True
r.font.color.rgb = NAVY

add_text(s5, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
         "a3brands.com  ·  Galaxy Visibility Diagnostic",
         size=13, color=DIM, font="Consolas", align=PP_ALIGN.CENTER)

if SAGGY_HEAD.exists():
    s5.shapes.add_picture(str(SAGGY_HEAD), Inches(11), Inches(0.3), height=Inches(1.4))

add_footer(s5)


prs.save(OUT)
print(f"wrote {OUT}")
